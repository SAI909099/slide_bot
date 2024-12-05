import asyncio
import sys
from os import getenv
import openai
from aiogram import Bot, Dispatcher, types, F
from aiogram.client.default import DefaultBotProperties
from aiogram.dispatcher import router
from aiogram.enums import ParseMode
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.types import Message
from pptx import Presentation
import logging

TOKEN = "7204615655:AAGGzAwbwIt4lQwvFaCmzhTi7z4J1hgMU20"
import torch
print(torch.__version__)
import tensorflow as tf
print(tf.__version__)

# bot = Bot(token=TOKEN)
# dp = Dispatcher()
bot = Bot(token=TOKEN, retry_after=3)  # Increase retry_after
dp = Dispatcher(storage=MemoryStorage())


@dp.message(F.text == "/start")
async def start(message: Message):
    await message.reply("Welcome! Send me a topic to start creating your presentation.")


def create_presentation(topic, slides_count):
    presentation = Presentation()
    for i in range(slides_count):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"{topic} - Slide {i+1}"
        slide.placeholders[1].text = generate_slide_content(topic, i+1)
    filename = f"presentations/{topic.replace(' ', '_')}.pptx"
    presentation.save(filename)
    return filename
# print(create_presentation("Artificial Intelligence", 1))

@dp.message()
async def handle_topic(message: Message):
    topic = message.text
    filename = create_presentation(topic, 5)  # Default to 5 slides
    await message.reply_document(types.FSInputFile(filename), caption="Here is your presentation!")



# openai.api_key = "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"


from transformers import pipeline

def generate_slide_content(topic, slide_number):
    generator = pipeline("text-generation", model="gpt2", framework="pt") # Replace with your preferred model
    prompt = f"Slide {slide_number} on the topic '{topic}':"
    response = generator(prompt, max_length=100, num_return_sequences=1)
    return response[0]['generated_text']
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

def create_pdf(topic, slides_content):
    filename = f"presentations/{topic.replace(' ', '_')}.pdf"
    pdf = canvas.Canvas(filename, pagesize=letter)
    for i, content in enumerate(slides_content, start=1):
        pdf.drawString(100, 750, f"Slide {i}: {content}")
        pdf.showPage()
    pdf.save()
    return filename


@dp.message()
async def handle_topic(message: Message):
    topic = message.text
    slides_content = [generate_slide_content(topic, i) for i in range(1, 6)]  # 5 slides
    ppt_filename = create_presentation(topic, 5)
    pdf_filename = create_pdf(topic, slides_content)

    await message.reply("Presentation generated! Choose a format:")
    await message.reply_document(types.FSInputFile(ppt_filename), caption="PPT Version")
    await message.reply_document(types.FSInputFile(pdf_filename), caption="PDF Version")




async def main() -> None:
    bot = Bot(token=TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.HTML))
    await dp.start_polling(bot)

if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, stream=sys.stdout)
    asyncio.run(main())
